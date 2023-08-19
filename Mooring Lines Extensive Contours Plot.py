# -*- coding: utf-8 -*-
"""
Created on Mon May 29 15:38:02 2023

@author: vince
"""

# %% Import packages
import moorpy as mp
import numpy as np
import pickle
from scipy.optimize import minimize, Bounds, NonlinearConstraint
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d.art3d import Poly3DCollection
from moorpy.MoorProps import getLineProps

from pptx.util import Inches


import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt


#  %% Import Libraries 

import numpy as np
import matplotlib.pyplot as plt
import os

#  %% Import functions

from Function_moorings import system,Thrust_calc,Thrust_IEA_15_MW_yaw,Thrust_IEA_15_MW,IEA_15_MW,\
                              Line_tension,transform_coordinates,turbine_geometry, \
                              turbine_geometry,turbine_2d_plot_xy,turbine_2d_plot_xz,turbine_2d_plot_yz,\
                              turbine_3d_plot,Thrust_plot_yaw,Thrust_plot_u,\
                              distance_anchor_fairlead, line_heading, slack_line_length, min_line_length, run_optimization, line_equation, movable_range,\
                              normal_iea_15mw_floater,Mooring_line_tension, Mooring_lines_optmization,\
                              Mooring_line_tension_ms, run_optimization_thrust_force_trial, save_simulation
                             
from pptx import Presentation
from pptx.util import Inches

# %% Create a presentation


from pptx.util import Inches

# %%
    
from pptx import Presentation
from PIL import Image


from pptx import Presentation
from pptx.util import Cm

def add_image_to_slide(presentation, slide_number, position, figure_path):
    # Check if the slide number is valid
    if slide_number < len(presentation.slides):
        slide = presentation.slides[slide_number]
    else:
        # The slide number is not valid, create a new slide
        slide_layout = presentation.slide_layouts[slide_number % len(presentation.slide_layouts)]
        slide = presentation.slides.add_slide(slide_layout)

    # Calculate maximum image dimensions on the slide
    slide_width = Cm(24)  # Width of the slide
    slide_height = Cm(18)  # Height of the slide
    space = Cm(0.05)  # Space between images
    max_image_width = (slide_width - space) / 2
    max_image_height = (slide_height - space) / 2

    # Load the image to get its original dimensions
    image = Image.open(figure_path)
    image_width, image_height = image.size

    # Calculate the scaling factor based on the maximum dimensions
    width_scale = max_image_width / image_width
    height_scale = max_image_height / image_height

    # Choose the smaller scaling factor to maintain the image's original proportions
    scale = min(width_scale, height_scale)

    # Calculate the new dimensions of the image
    new_image_width = image_width * scale
    new_image_height = image_height * scale

    # Calculate image position based on the specified position
    left = (slide_width - 2 * max_image_width - space + Cm(1) )/ 2 + position % 2 * (max_image_width + space)
    top = (slide_height - 2 * max_image_height - space + Cm(1) ) / 2 + position // 2 * (max_image_height + space)

    # Add the image to the slide
    slide.shapes.add_picture(figure_path, left, top, width=new_image_width, height=new_image_height)





# %% Set geometry parameters of the mooring system
depth       = 200       # Water depth [m]
zFair       = -14       # Fairlead z elevation [m]
rAnchor     = 837.60    # Anchor radius [m]
rFair       = 58        # Fairlead radius [m]
lineLength  = 850       # Line unstretched length [m]

# %% Definition of the movanle range domain

n_cont=30


# Obtain movable range vertex coordinates [m]
A,B,C = movable_range()

# Linear equations that describes the upper and lower boundary of movable range
y1 = line_equation(A,B)
y2 = line_equation(A,C)

# Create a list of x- and y- coordinates [m]    
x_list = np.linspace(A[0], B[0], n_cont)
y_list = np.linspace(C[1], B[1], n_cont)

# Creating 2-D grid
[X, Y] = np.meshgrid(x_list, y_list)

# %% The Thust force is conìomputed and applyies to the system

ws_list= [11]
wd_list= [0]
yaw_list = [0]

merged_presentation= Presentation()

for ws in ws_list:
    for wd in wd_list:
        for yaw in yaw_list:
            
            presentation = Presentation()

            # Create the presentation name with variable values
            presentation_name = f'Contours at ws={ws} wd={wd} yaw={yaw}.pptx'
            
            # Add the title slide with the name of the presentation
            slide_layout = presentation.slide_layouts[0]  # Title slide layout
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f'Contours at ws={ws} m/s, wd={wd}°, yaw={yaw}°'
            
            
            #%%
            path_to_folder = 'Optimization Results/'+'%0.1f ws' %ws +'/' + '%i wd' %wd +'/' + '%i yaw' %yaw +'/'
            
                
            with open(path_to_folder + 'optimization.pickle', 'rb') as f:
                  optimization = pickle.load(f)
                  
                  line_length1 =optimization['line_length1']
                  line_length2 =optimization['line_length2']
                  line_length3 =optimization['line_length3']
            
                  error_x =optimization['error_x']
                  error_y =optimization['error_y']
                  error_dist =optimization['error_dist']
                  error_psi=optimization['error_psi']
            
                  
                  TF1 =optimization['TF1']
                  TF2 =optimization['TF2']
                  TF3 =optimization['TF3']
                  HF1 =optimization['HF1']
                  HF2 =optimization['HF2']
                  HF3 =optimization['HF3']
                  VF1 =optimization['VF1']
                  VF2 =optimization['VF2']
                  VF3 =optimization['VF3']
                  
                  K11 =optimization['K11']
                  K12 =optimization['K12']
                  K13 =optimization['K13']
                  K21 =optimization['K21']
                  K22 =optimization['K22']
                  K23 =optimization['K23']
                  K31 =optimization['K31']
                  K32 =optimization['K32']
                  K33 =optimization['K33']
            
            
            #%%
              
            path_to_folder_contours = 'Optimization Results/'+'%0.1f ws' %ws +'/' + '%i wd' %wd +'/' + '%i yaw' %yaw +'/Contours/'
            
            # Create the folder if it doesn't exist
            if not os.path.exists(path_to_folder_contours):
                os.makedirs(path_to_folder_contours)
            
            
            dpi_val = 600
            
            # %% Distance between the desired and actual position
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, np.sqrt((error_x)**2+(error_y)**2),np.arange(0, 100, 10),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.minorticks_on()
            cb.set_label('Distance [m]')
            title='Distance between the desired and actual position'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 0
            add_image_to_slide(presentation, 1, 0, figure_path)
            
            # %% Contour plot for the absolute error in yaw angle of the floater
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, np.abs(np.rad2deg(error_psi)),np.arange(0, 15,1),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.minorticks_on()
            cb.set_label('Error [deg]')
            title='Absolute error in yaw angle of the floater'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 1, 1, figure_path)
            
# %%
            def contour_plot(X,Y,Z,label,title,x,y,slide,position)  :
            
                fig, ax = plt.subplots(1, 1)
                cp = ax.contourf(X, Y, Z)
                # cp = ax.contourf(X, Y, Z,np.arange(0, 1600, 200),
                #                     extend='both')
                cb = fig.colorbar(cp)
                cb.set_label(label)
                cb.minorticks_on()
                ax.set_title(title)
                ax.set_xlabel(x)
                ax.set_ylabel(y)
                ax.axis('equal')
                plt.grid()
                fig.set_dpi(dpi_val)
                plt.show()
            
                
                # Save the plot with the title as the figure name
                figure_name = f'{title.replace(" ", "_")}.png'
                figure_path = os.path.join(path_to_folder_contours, figure_name)
                plt.savefig(figure_path,bbox_inches='tight',  dpi=600)
            
                # Add the image to the slide                
                add_image_to_slide(presentation, slide, position, figure_path)
            
            
                return 
               
            #sum of the lines length
            contour_plot(X,Y,line_length1+line_length2+line_length3,'Line length [m]','Sum of the line lengths','x [m]','y [m]',2,0) 
            
            #line length 1 
            contour_plot(X,Y,line_length1,'Line length [m]','Line length 1','x [m]','y [m]',2,1) 
            
            #line length 2 
            contour_plot(X,Y,line_length2,'Line length [m]','Line length 2','x [m]','y [m]',2,2) 
            
            #line length 3 
            contour_plot(X,Y,line_length3,'Line length [m]','Line length 3','x [m]','y [m]',2,3) 
            
            

            # %% Sum of the tensions
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (TF1+TF2+TF3)/10**6,np.arange(6, 15, 0.9),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Tension [MN]')
            cb.minorticks_on()
            title='Sum of the tensions'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 3, 0, figure_path)
            
            # %% Tension line 1
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (TF1)/10**6,np.arange(1, 10, 1),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Tension [MN]')
            cb.minorticks_on()
            title='Tension line 1'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 3, 1, figure_path)
            
            # %% Tension line 2
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (TF2)/10**6,np.arange(1, 10, 1),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Tension [MN]')
            cb.minorticks_on()
            title='Tension line 2'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 3, 2, figure_path)
            
            # %% Tension line 3
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (TF3)/10**6,np.arange(1, 10, 1),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Tension [MN]')
            cb.minorticks_on()
            title='Tension line 3'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 3, 3, figure_path)
            # %%Sum of horizontal tensions
                
            # Plot the figure
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (HF1+HF2+HF3)/10**6, np.arange(3, 13, 0.9), extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Horizontal tension [MN]')
            cb.minorticks_on()
            
            title = 'Sum of the horizontal tensions'
            ax.set_title(title)
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path,bbox_inches='tight',  dpi=600)
            
            # Add the image to the slide
            add_image_to_slide(presentation, 4, 0, figure_path)
            
            
            # %% Horizontal tension 1
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (HF1)/10**6,np.arange(0, 9, 0.9),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Horizontal tension [MN]')
            cb.minorticks_on()
            
            title = 'Horizontal tension line 1'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to the slide
            add_image_to_slide(presentation, 4, 1, figure_path)
            
            # %% Horizontal tension 2
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (HF2)/10**6,np.arange(0, 9, 0.9),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Horizontal tension [MN]')
            cb.minorticks_on()
            
            title = 'Horizontal tension line 2'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to the slide
            add_image_to_slide(presentation, 4, 2, figure_path)
            
            # %% Horizontal tension 3
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (HF3)/10**6,np.arange(0, 9, 0.9),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Horizontal tension [MN]')
            cb.minorticks_on()
            
            title = 'Horizontal tension line 3'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to the slide
            add_image_to_slide(presentation, 4, 3, figure_path)
            
            # %% Sum of the vertical tensions
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (VF1+VF2+VF3)/10**6,np.arange(5, 10, 0.5),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Vertical tension [MN]')
            cb.minorticks_on()
            
            title='Sum of the vertical tensions'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 5, 0, figure_path)
            
            # %% Vertical tension 1
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (VF1)/10**6,np.arange(1, 5, 0.5),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Vertical tension [MN]')
            cb.minorticks_on()
            title='Vertical tension line 1'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 5, 1, figure_path)
            
            # %% Vertical tension 2
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (VF2)/10**6,np.arange(1, 5, 0.5),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Vertical tension [MN]')
            cb.minorticks_on()
            title='Vertical tension line 2'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 5, 2, figure_path)
            
            # %% Vertical tension 3
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, (VF3)/10**6,np.arange(1, 5, 0.5),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.set_label('Vertical tension [MN]')
            cb.minorticks_on()
            title='Vertical tension line 3'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 5, 3, figure_path)
            
            
            
            # %%
            # Contour plot for the difference between the stiffness of the mooring system at 
            # (0,0) and at the desidered position for the 
            # for the  first degree of freedom(sway) of the mooring system
            
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, K11/10**3,np.arange(0, 500,50),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.minorticks_on()
            cb.set_label(' K11 [kN/m]')
            title='Mooring lines system stiffness Surge term K11'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 6, 0, figure_path)
            # %%
            
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, K22/10**3,np.arange(0, 500,50),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.minorticks_on()
            cb.set_label(' K22 [kN/m]')
            title='Mooring lines system stiffness Sway term K22'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 6, 1, figure_path)
            # %%
            
            fig, ax = plt.subplots(1, 1)
            cp = ax.contourf(X, Y, K33/10**3,np.arange(0, 500,50),
                                extend='both')
            cb = fig.colorbar(cp)
            cb.minorticks_on()
            cb.set_label(' K33 [kNm/rad]')
            title='Mooring lines system stiffness Yaw term K33'
            ax.set_title(title)
            
            ax.set_xlabel('x [m]')
            ax.set_ylabel('y [m]')
            ax.axis('equal')
            plt.grid()
            fig.set_dpi(dpi_val)
            plt.show()
            
            # Save the plot with the title as the figure name
            figure_name = f'{title.replace(" ", "_")}.png'
            figure_path = os.path.join(path_to_folder_contours, figure_name)
            plt.savefig(figure_path, bbox_inches='tight', dpi=600)
            
            # Add the image to slide number 2
            add_image_to_slide(presentation, 6, 3, figure_path)
            # %%
            
            # # Check if the presentation file already exists
            # if os.path.exists(presentation_name):
            #     # Find the next available number for the presentation name
            #     index = 1
            #     while os.path.exists(f'{presentation_name}_{index}.pptx'):
            #         index += 1
            #     # Add the number to the presentation name
            #     presentation_name = f'{presentation_name}_{index}.pptx'
            
            # Save the presentation
            presentation.save(presentation_name)
            
            print('Saved:',presentation_name)
            
            
#             for slide in presentation.slides:
#                 # Append each slide to the merged presentation
#                 merged_presentation.slides.add_slide(slide)
            
# # Save the merged presentation
# merged_presentation.save('All_contours.pptx')